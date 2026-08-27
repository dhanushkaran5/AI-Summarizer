"""Comprehensive document text extraction service supporting PDF, DOCX, PPTX, TXT, MD, HTML, and CSV."""
import os
import re
import csv
from typing import Optional


def extract_text(file_path: str, file_type: str) -> dict:
    """Extract text and structural pages/sections from any supported document format.
    
    Returns:
        dict with:
            text: full extracted string
            pages: list of {page_number, text} dicts
            page_count: total pages/divisions
            has_scanned_content: boolean flag for OCR recommendation
            structure: list of detected headings
    """
    extractors = {
        "pdf": _extract_pdf,
        "docx": _extract_docx,
        "pptx": _extract_pptx,
        "txt": _extract_txt,
        "md": _extract_markdown,
        "markdown": _extract_markdown,
        "html": _extract_html,
        "htm": _extract_html,
        "csv": _extract_csv,
    }

    normalized_type = file_type.lower().replace(".", "")
    extractor = extractors.get(normalized_type)
    if not extractor:
        raise ValueError(f"Unsupported file type: {file_type}. Supported: PDF, DOCX, PPTX, TXT, MD, HTML, CSV")

    return extractor(file_path)


def _extract_pdf(file_path: str) -> dict:
    """Extract text from PDF with page tracking and scanned-page OCR detection."""
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(file_path)
        pages = []
        full_text = []
        low_text_page_count = 0

        for i, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text() or ""
            except Exception:
                page_text = ""
            
            page_text = page_text.strip()
            if len(page_text) < 30:  # Very little or no text might indicate scanned/image page
                low_text_page_count += 1

            if page_text:
                pages.append({"page_number": i + 1, "text": page_text})
                full_text.append(page_text)

        total_pages = len(reader.pages)
        is_scanned = (total_pages > 0 and (low_text_page_count / total_pages) > 0.6) or (len(full_text) == 0)

        extracted_text = "\n\n".join(full_text) if full_text else ""

        return {
            "text": extracted_text,
            "pages": pages,
            "page_count": max(1, total_pages),
            "has_scanned_content": is_scanned,
            "ocr_recommended": is_scanned and len(extracted_text) < 100,
        }
    except Exception as e:
        raise ValueError(f"Failed to extract PDF: {str(e)}")


def _extract_docx(file_path: str) -> dict:
    """Extract text from DOCX with paragraph and section tracking."""
    try:
        from docx import Document

        doc = Document(file_path)
        full_text = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                full_text.append(text)

        all_text = "\n\n".join(full_text)
        word_count = len(all_text.split())
        estimated_pages = max(1, word_count // 250)

        pages = []
        if full_text:
            chunk_size = max(1, len(full_text) // estimated_pages)
            for i in range(0, len(full_text), chunk_size):
                page_text = "\n\n".join(full_text[i:i + chunk_size])
                pages.append({
                    "page_number": len(pages) + 1,
                    "text": page_text,
                })

        return {
            "text": all_text,
            "pages": pages if pages else [{"page_number": 1, "text": all_text}],
            "page_count": len(pages) if pages else 1,
            "has_scanned_content": False,
            "ocr_recommended": False,
        }
    except Exception as e:
        raise ValueError(f"Failed to extract DOCX: {str(e)}")


def _extract_pptx(file_path: str) -> dict:
    """Extract text from PPTX with slide-level tracking."""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        pages = []
        full_text = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text.append(text)

            if slide_text:
                combined_slide = "\n".join(slide_text)
                pages.append({"page_number": i + 1, "text": combined_slide})
                full_text.append(f"## Slide {i+1}\n{combined_slide}")

        return {
            "text": "\n\n".join(full_text),
            "pages": pages if pages else [{"page_number": 1, "text": "Empty presentation"}],
            "page_count": max(1, len(prs.slides)),
            "has_scanned_content": False,
            "ocr_recommended": False,
        }
    except Exception as e:
        raise ValueError(f"Failed to extract PPTX: {str(e)}")


def _extract_txt(file_path: str) -> dict:
    """Extract text from TXT file with encoding detection."""
    try:
        with open(file_path, "rb") as f:
            raw_data = f.read()

        try:
            import chardet
            detected = chardet.detect(raw_data)
            encoding = detected.get("encoding", "utf-8") or "utf-8"
        except Exception:
            encoding = "utf-8"

        text = raw_data.decode(encoding, errors="replace").strip()

        words = text.split()
        pages = []
        page_size = 250
        for i in range(0, max(1, len(words)), page_size):
            page_text = " ".join(words[i:i + page_size])
            if page_text:
                pages.append({"page_number": len(pages) + 1, "text": page_text})

        return {
            "text": text,
            "pages": pages if pages else [{"page_number": 1, "text": text}],
            "page_count": max(1, len(pages)),
            "has_scanned_content": False,
            "ocr_recommended": False,
        }
    except Exception as e:
        raise ValueError(f"Failed to extract TXT: {str(e)}")


def _extract_markdown(file_path: str) -> dict:
    """Extract and structure Markdown document."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read().strip()

        # Split by top-level headers (# or ##) to create natural sections
        sections = re.split(r'\n(?=#{1,2}\s)', text)
        pages = []
        for i, s in enumerate(sections):
            if s.strip():
                pages.append({"page_number": i + 1, "text": s.strip()})

        return {
            "text": text,
            "pages": pages if pages else [{"page_number": 1, "text": text}],
            "page_count": max(1, len(pages)),
            "has_scanned_content": False,
            "ocr_recommended": False,
        }
    except Exception as e:
        raise ValueError(f"Failed to extract Markdown: {str(e)}")


def _extract_html(file_path: str) -> dict:
    """Extract clean text and structure from HTML."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            html_content = f.read()

        # Strip scripts, styles, and tags
        clean_html = re.sub(r'<script.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
        clean_html = re.sub(r'<style.*?</style>', '', clean_html, flags=re.DOTALL | re.IGNORECASE)
        
        # Replace line breaks and headers with newlines
        clean_html = re.sub(r'<(?:p|div|h[1-6]|li|tr)[^>]*>', '\n', clean_html, flags=re.IGNORECASE)
        clean_text = re.sub(r'<[^>]+>', ' ', clean_html)
        clean_text = re.sub(r'[ \t]+', ' ', clean_text)
        clean_text = re.sub(r'\n\s*\n+', '\n\n', clean_text).strip()

        words = clean_text.split()
        pages = []
        page_size = 250
        for i in range(0, max(1, len(words)), page_size):
            page_text = " ".join(words[i:i + page_size])
            if page_text:
                pages.append({"page_number": len(pages) + 1, "text": page_text})

        return {
            "text": clean_text,
            "pages": pages if pages else [{"page_number": 1, "text": clean_text}],
            "page_count": max(1, len(pages)),
            "has_scanned_content": False,
            "ocr_recommended": False,
        }
    except Exception as e:
        raise ValueError(f"Failed to extract HTML: {str(e)}")


def _extract_csv(file_path: str) -> dict:
    """Extract and summarize structured tabular CSV data."""
    try:
        rows = []
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                if any(cell.strip() for cell in row):
                    rows.append(row)

        if not rows:
            return {"text": "Empty CSV dataset", "pages": [{"page_number": 1, "text": "Empty CSV"}], "page_count": 1, "has_scanned_content": False, "ocr_recommended": False}

        header = rows[0]
        header_text = " | ".join(header)
        text_lines = [f"Columns: {header_text}\nTotal Records: {len(rows) - 1}\n"]

        pages = []
        chunk_size = 30
        for page_idx, i in enumerate(range(1, len(rows), chunk_size)):
            batch = rows[i:i + chunk_size]
            batch_lines = [f"Row {i + j}: " + ", ".join(f"{header[k] if k < len(header) else f'Col{k}'}: {val}" for k, val in enumerate(r)) for j, r in enumerate(batch)]
            page_content = "\n".join(batch_lines)
            pages.append({"page_number": page_idx + 1, "text": page_content})
            text_lines.append(page_content)

        full_text = "\n\n".join(text_lines)
        return {
            "text": full_text,
            "pages": pages if pages else [{"page_number": 1, "text": full_text}],
            "page_count": max(1, len(pages)),
            "has_scanned_content": False,
            "ocr_recommended": False,
        }
    except Exception as e:
        raise ValueError(f"Failed to extract CSV: {str(e)}")
